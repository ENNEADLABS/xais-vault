"""
PPTX extractor — uses python-pptx for PowerPoint text extraction.
Extracts slide titles and text shapes, preserving slide boundaries.
"""

import logging

from pptx import Presentation

from . import ExtractionResult

logger = logging.getLogger(__name__)


async def extract_pptx(file_path: str) -> ExtractionResult:
    """Extract text from a PPTX file, preserving slide structure."""
    prs = Presentation(file_path)
    slides: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        parts: list[str] = [f"--- Slide {slide_num} ---"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)

            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append(" | ".join(cells))

        # Only include slides with actual content
        if len(parts) > 1:
            slides.append("\n".join(parts))

    full_text = "\n\n".join(slides)
    word_count = len(full_text.split())

    logger.info(f"PPTX extracted: {len(slides)} slides, {word_count} words — {file_path}")

    return ExtractionResult(
        text=full_text,
        page_count=len(slides),
        word_count=word_count,
        metadata={"extractor": "python-pptx", "slide_count": len(slides)},
    )
